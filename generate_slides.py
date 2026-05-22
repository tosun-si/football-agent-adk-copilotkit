"""Generate the talk .pptx from slides_content.py.

Usage:
    uv sync
    uv run python generate_slides.py

Output: slides/talk_devlille_gcs.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from slides_content import SLIDES

ROOT = Path(__file__).parent
OUTPUT = ROOT / "slides" / "talk_devlille_gcs.pptx"

# Palette
NAVY = RGBColor(0x1A, 0x2B, 0x4A)
DARK_GRAY = RGBColor(0x2C, 0x2C, 0x2C)
ORANGE = RGBColor(0xE8, 0x59, 0x0C)
LIGHT_GRAY = RGBColor(0xEE, 0xEE, 0xEE)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)
CODE_FG = RGBColor(0xE6, 0xE6, 0xE6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_run(run, *, size=None, bold=None, italic=None, color=None, font=None):
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    if font is not None:
        run.font.name = font


def _add_title(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run(run, size=36, bold=True, color=NAVY)
    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(1.2), Inches(0.07)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def _add_notes(slide, notes: str | None):
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _add_image_or_placeholder(slide, image_path: str, placeholder: str = ""):
    abs_path = Path(image_path)
    if not abs_path.is_absolute():
        abs_path = ROOT / image_path

    body_top = Inches(1.5)
    body_left = Inches(1)
    body_width = Inches(11.3)
    body_height = Inches(5.2)

    if image_path and abs_path.exists():
        # Fit by width, let height adjust automatically
        slide.shapes.add_picture(
            str(abs_path), body_left, body_top, width=body_width
        )
        return

    # Placeholder rectangle
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, body_left, body_top, body_width, body_height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT_GRAY
    rect.line.color.rgb = DARK_GRAY
    rect.line.dash_style = 7  # dashed
    rect.line.width = Pt(1.5)

    tf = rect.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ {placeholder or 'Image à insérer'} ]"
    _set_run(run, size=18, italic=True, color=DARK_GRAY)


def _add_caption(slide, caption: str):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = caption
    _set_run(run, size=14, italic=True, color=DARK_GRAY)


# === Renderers ===

def render_title(prs, data):
    slide = _blank(prs)
    # Big centered title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = data["title"]
    _set_run(run, size=54, bold=True, color=NAVY)

    if data.get("subtitle"):
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.3), Inches(1.2))
        tf = sub.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = data["subtitle"]
        _set_run(run, size=24, color=DARK_GRAY)

    # Accent bar under subtitle
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.66), Inches(5.7), Inches(2), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    if data.get("footer"):
        foot = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6))
        tf = foot.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = data["footer"]
        _set_run(run, size=16, color=ORANGE)

    _add_notes(slide, data.get("notes"))


def render_bullets(prs, data):
    slide = _blank(prs)
    _add_title(slide, data["title"])
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(data["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"•   {bullet}"
        _set_run(run, size=22, color=DARK_GRAY)
        p.space_after = Pt(16)
    _add_notes(slide, data.get("notes"))


def render_image(prs, data):
    slide = _blank(prs)
    _add_title(slide, data["title"])
    _add_image_or_placeholder(
        slide, data.get("image_path", ""), data.get("image_placeholder", "")
    )
    if data.get("caption"):
        _add_caption(slide, data["caption"])
    _add_notes(slide, data.get("notes"))


def render_code(prs, data):
    slide = _blank(prs)
    _add_title(slide, data["title"])

    # Dark code block
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.5)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = CODE_BG
    rect.line.fill.background()

    tf = rect.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)

    code_lines = data["code"].split("\n")
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line if line else " "
        _set_run(run, size=14, color=CODE_FG, font="Consolas")
        p.space_after = Pt(0)

    _add_notes(slide, data.get("notes"))


def render_demo(prs, data):
    slide = _blank(prs)
    # Black background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLACK
    bg.line.fill.background()

    # Big DEMO
    demo = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.3), Inches(2))
    tf = demo.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "DEMO"
    _set_run(run, size=140, bold=True, color=ORANGE)

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.8))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = data["title"]
    _set_run(run, size=28, color=WHITE)

    # Steps
    if data.get("steps"):
        steps = slide.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(8.3), Inches(2.5))
        tf = steps.text_frame
        tf.word_wrap = True
        for i, step in enumerate(data["steps"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = f"{i + 1}.   {step}"
            _set_run(run, size=20, color=WHITE)
            p.space_after = Pt(10)

    _add_notes(slide, data.get("notes"))


def render_takeaways(prs, data):
    slide = _blank(prs)
    _add_title(slide, data["title"])

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True

    for i, (headline, sub) in enumerate(data["items"]):
        # Headline
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"{i + 1}.   {headline}"
        _set_run(run, size=28, bold=True, color=NAVY)
        p.space_after = Pt(6)
        # Sub
        sub_p = tf.add_paragraph()
        sub_run = sub_p.add_run()
        sub_run.text = f"        {sub}"
        _set_run(sub_run, size=18, color=DARK_GRAY)
        sub_p.space_after = Pt(28)

    _add_notes(slide, data.get("notes"))


RENDERERS = {
    "title": render_title,
    "bullets": render_bullets,
    "image": render_image,
    "code": render_code,
    "demo": render_demo,
    "takeaways": render_takeaways,
}


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for i, slide_data in enumerate(SLIDES, start=1):
        renderer = RENDERERS.get(slide_data["type"])
        if renderer is None:
            raise ValueError(f"Slide {i}: unknown type {slide_data['type']!r}")
        renderer(prs, slide_data)
    return prs


def main():
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(OUTPUT)
    print(f"OK — {len(prs.slides)} slides → {OUTPUT}")


if __name__ == "__main__":
    main()
